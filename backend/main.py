import os
import io
import json
import re
import base64
from pathlib import Path
from functools import lru_cache

import anthropic
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import Literal
from supabase import create_client, Client

app = FastAPI(title="Mapa Mental - Enfermagem")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@lru_cache(maxsize=1)
def get_supabase() -> Client:
    url = os.environ.get("SUPABASE_URL", "")
    key = os.environ.get("SUPABASE_KEY", "")
    if not url or not key:
        raise RuntimeError("SUPABASE_URL ou SUPABASE_KEY não configuradas.")
    return create_client(url, key)


SYSTEM_PROMPT = """Você é um assistente especializado em criar materiais de estudo para estudantes de enfermagem.
Analise o conteúdo de PDFs de aulas e slides e produza:

1. Um mapa mental como árvore JSON com categorias visuais
2. Um resumo estruturado com os principais tópicos

CATEGORIAS DISPONÍVEIS para cada nó:
- "root": tópico central (apenas o nó raiz)
- "definition": definição, conceito, introdução
- "pathophysiology": fisiopatologia, mecanismo, etiologia, causas
- "symptoms": sinais, sintomas, manifestações clínicas
- "diagnosis": diagnóstico, exames, critérios diagnósticos
- "treatment": tratamento, medicamentos, conduta terapêutica
- "nursing": cuidados de enfermagem, SAE, intervenções de enfermagem
- "classification": classificação, tipos, formas, estágios
- "epidemiology": epidemiologia, prevalência, fatores de risco
- "detail": detalhes, sub-itens, informações complementares

REGRAS DO MAPA MENTAL:
- Máximo 4 níveis de profundidade
- IDs únicos: "root", "n1", "n1-1", "n1-2", "n2", "n2-1", etc.
- Labels concisos (máx 8 palavras por nó)
- Atribua a categoria mais específica possível a cada nó
- Inclua todos os temas principais do material

REGRAS DO RESUMO:
- Identifique o tópico principal
- Liste de 5 a 10 pontos-chave clínicos
- Organize por seções temáticas
- Foque nos aspectos mais cobrados em provas

Retorne EXCLUSIVAMENTE um JSON válido no formato:
{
  "mindmap": {
    "id": "root",
    "label": "Nome do Tópico Principal",
    "category": "root",
    "children": [
      {
        "id": "n1",
        "label": "Nome do Ramo",
        "category": "definition",
        "children": [
          { "id": "n1-1", "label": "Detalhe", "category": "detail", "children": [] }
        ]
      }
    ]
  },
  "summary": {
    "main_topic": "Nome do tópico principal",
    "key_points": ["ponto 1", "ponto 2"],
    "sections": [
      { "title": "Título da seção", "points": ["ponto A", "ponto B"] }
    ]
  }
}"""

MAX_FILE_BYTES = 30 * 1024 * 1024

ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt",
    ".jpg", ".jpeg", ".png", ".gif", ".webp",
    ".txt", ".md",
}

IMAGE_MEDIA_TYPES = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".gif": "image/gif",
    ".webp": "image/webp",
}


def extract_docx_text(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def build_content_block(filename: str, ext: str, data: bytes) -> list:
    """Converte um arquivo em um ou mais blocos de conteúdo para a API do Claude."""
    b64 = base64.standard_b64encode(data).decode("utf-8")

    if ext == ".pdf":
        return [
            {"type": "text", "text": f"Arquivo: {filename}"},
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": b64}},
        ]

    if ext in IMAGE_MEDIA_TYPES:
        return [
            {"type": "text", "text": f"Imagem: {filename}"},
            {"type": "image", "source": {"type": "base64", "media_type": IMAGE_MEDIA_TYPES[ext], "data": b64}},
        ]

    if ext in (".docx", ".doc"):
        text = extract_docx_text(data)
        return [{"type": "text", "text": f"Arquivo Word ({filename}):\n{text}"}]

    if ext in (".pptx", ".ppt"):
        text = extract_pptx_text(data)
        return [{"type": "text", "text": f"Apresentação PowerPoint ({filename}):\n{text}"}]

    if ext in (".txt", ".md"):
        text = data.decode("utf-8", errors="replace")
        return [{"type": "text", "text": f"Arquivo de texto ({filename}):\n{text}"}]

    return []


class ProcessResponse(BaseModel):
    id: str
    mindmap: dict
    summary: dict
    diagram: str
    files_processed: list[str]
    created_at: str


def _build_diagram_content(summary: dict) -> str:
    sections = summary.get("sections", [])
    sections_text = "\n".join(
        f'{s["title"]}:\n' + "\n".join(f'  - {p}' for p in s["points"])
        for s in sections
    )
    return (
        f'Tópico: {summary["main_topic"]}\n\n'
        f'Pontos-chave:\n' + "\n".join(f'- {p}' for p in summary.get("key_points", [])) +
        f'\n\nSeções:\n{sections_text}\n\n'
        f'Gere o infográfico SVG educacional e aprofundado conforme as instruções. '
        f'Para cada categoria relevante (sintomas, tratamento, fisiopatologia, procedimentos, teorias, etc.), '
        f'escreva EXPLICAÇÕES COMPLETAS de 2 a 4 frases — não apenas tópicos curtos. '
        f'O objetivo é que uma estudante de enfermagem consiga estudar apenas pelo infográfico.'
    )


def _generate_diagram_svg(client: anthropic.Anthropic, summary: dict) -> str:
    """Gera o SVG do diagrama a partir do resumo. Retorna string vazia em caso de falha."""
    try:
        resp = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=8192,
            system=SVG_SYSTEM,
            messages=[{"role": "user", "content": _build_diagram_content(summary)}],
        )
        svg = resp.content[0].text.strip()
        if svg.startswith("```"):
            svg = re.sub(r"^```(?:svg|xml)?\s*", "", svg)
            svg = re.sub(r"\s*```$", "", svg.strip())
        return svg
    except Exception:
        return ""


@app.post("/api/process", response_model=ProcessResponse)
async def process_pdfs(files: list[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="Nenhum arquivo enviado.")

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY não configurada.")

    client = anthropic.Anthropic(api_key=api_key)
    content_blocks: list = []
    processed_names = []

    for upload in files:
        ext = Path(upload.filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"Formato '{ext}' não suportado. Use: PDF, Word, PowerPoint, imagens (JPG/PNG/WEBP) ou TXT.",
            )

        file_bytes = await upload.read()
        if len(file_bytes) > MAX_FILE_BYTES:
            raise HTTPException(status_code=413, detail=f"'{upload.filename}' é muito grande (máx 30 MB).")

        blocks = build_content_block(upload.filename, ext, file_bytes)
        if not blocks:
            raise HTTPException(status_code=400, detail=f"Não foi possível processar '{upload.filename}'.")

        content_blocks.extend(blocks)
        processed_names.append(upload.filename)

    content_blocks.append({
        "type": "text",
        "text": "Analise todos os documentos acima e gere o mapa mental e o resumo conforme instruído. Retorne apenas o JSON, sem texto antes ou depois."
    })

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=8192,
            system=SYSTEM_PROMPT,
            messages=[{"role": "user", "content": content_blocks}],
        )
    except anthropic.APIError as e:
        raise HTTPException(status_code=502, detail=f"Erro na API do Claude: {str(e)}")

    raw = response.content[0].text.strip()

    # Remove markdown code fences se presentes
    if "```" in raw:
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```\s*$", "", raw)

    # Extrai o JSON mesmo que haja texto antes/depois
    match = re.search(r'\{[\s\S]*\}', raw)
    if match:
        raw = match.group(0)

    try:
        result = json.loads(raw)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=502, detail=f"Resposta da IA não é um JSON válido: {str(e)[:120]}")

    # Salva no Supabase primeiro (sem diagrama) para garantir que o resultado seja persistido
    try:
        sb = get_supabase()
        payload: dict = {
            "topic": result["summary"]["main_topic"],
            "mindmap": json.dumps(result["mindmap"], ensure_ascii=False),
            "summary": result["summary"],
            "files_processed": processed_names,
        }
        row = sb.table("mindmaps").insert(payload).execute()
        saved = row.data[0]
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Erro ao salvar no banco: {str(e)}")

    # Gera diagrama SVG após salvar (falha silenciosa se der erro)
    diagram_svg = _generate_diagram_svg(client, result["summary"])

    # Atualiza o registro com o diagrama, se a coluna existir
    if diagram_svg:
        try:
            sb.table("mindmaps").update({"diagram": diagram_svg}).eq("id", saved["id"]).execute()
        except Exception:
            pass  # Coluna ainda não existe — sem problema, usuário pode gerar depois

    return ProcessResponse(
        id=saved["id"],
        mindmap=result["mindmap"],
        summary=result["summary"],
        diagram=diagram_svg,
        files_processed=processed_names,
        created_at=saved["created_at"],
    )


@app.get("/api/history")
def get_history():
    try:
        sb = get_supabase()
        rows = sb.table("mindmaps").select("*").order("created_at", desc=True).limit(30).execute()
        # Desserializa mindmap de string para dict
        for row in rows.data:
            if isinstance(row.get("mindmap"), str):
                try:
                    row["mindmap"] = json.loads(row["mindmap"])
                except Exception:
                    row["mindmap"] = {"id": "root", "label": row.get("topic", ""), "category": "root", "children": []}
            if not row.get("diagram"):
                row["diagram"] = ""
        return rows.data
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Erro ao carregar histórico: {str(e)}")


@app.delete("/api/history/{item_id}")
def delete_history(item_id: str):
    try:
        sb = get_supabase()
        sb.table("mindmaps").delete().eq("id", item_id).execute()
        return {"ok": True}
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Erro ao deletar: {str(e)}")


SVG_SYSTEM = """Você é um especialista em enfermagem e designer de materiais educativos em SVG.
Sua tarefa é gerar um infográfico SVG completo e APROFUNDADO que sirva como guia de estudo autossuficiente.

FILOSOFIA DO CONTEÚDO:
- Cada item deve ser uma EXPLICAÇÃO REAL de 1 a 3 frases — não palavras soltas
- Explique o MECANISMO, não só o nome (ex: "A febre ocorre porque..." em vez de só "Febre")
- Para tratamentos: explique a lógica terapêutica (ex: "Administra-se X porque atua em Y, promovendo Z")
- Para procedimentos: descreva as etapas e o racional clínico
- Para teorias: explique o conceito e sua aplicação prática na enfermagem
- Para sintomas: explique por que ocorrem (fisiopatologia do sintoma)
- O objetivo: uma estudante consegue estudar APENAS pelo infográfico, sem precisar do PDF original

ESTRUTURA DO SVG:

1. DIMENSÕES: width="960". Height calculado pelo conteúdo (mínimo 800).

2. FUNDO: <rect width="960" height="TOTAL" fill="#f4f6f9"/>

3. CABEÇALHO (y=0, height=90):
   <rect width="960" height="90" fill="#2d6a5f"/>
   Subtítulo (especialidade): <text x="480" y="30" text-anchor="middle" font-family="Arial,sans-serif" font-size="11" fill="rgba(255,255,255,0.75)" letter-spacing="2">ENFERMAGEM PEDIÁTRICA</text>
   Título principal: <text x="480" y="62" text-anchor="middle" font-family="Arial,sans-serif" font-size="22" font-weight="bold" fill="white">TÍTULO DO TÓPICO</text>

4. GRID: 2 colunas, x=20 (col1) e x=500 (col2), largura=440px cada, gap=20px.
   Primeira linha de cards: y=110.

5. ESTRUTURA DE CADA CARD:
   <!-- fundo do card -->
   <rect x="X" y="Y" width="440" height="ALTURA_TOTAL" rx="12" fill="COR_BG" stroke="COR_BORDA" stroke-width="1.5"/>
   <!-- header colorido -->
   <rect x="X" y="Y" width="440" height="46" rx="12" fill="COR_HEADER"/>
   <rect x="X" y="Y+24" width="440" height="22" fill="COR_HEADER"/>
   <!-- ícone redondo no header -->
   <circle cx="X+22" cy="Y+23" r="10" fill="rgba(255,255,255,0.2)"/>
   <text x="X+22" y="Y+28" text-anchor="middle" font-family="Arial" font-size="12" fill="white">●</text>
   <!-- título do card -->
   <text x="X+40" y="Y+29" font-family="Arial,sans-serif" font-size="13" font-weight="bold" fill="white">TÍTULO DO CARD</text>
   <!-- corpo com texto aprofundado via foreignObject -->
   <foreignObject x="X" y="Y+46" width="440" height="ALTURA_CORPO">
     <div xmlns="http://www.w3.org/1999/xhtml" style="font-family:Arial,sans-serif;font-size:12px;color:#1e1e2e;padding:14px 16px 10px;line-height:1.65">
       <div style="margin-bottom:10px;padding-left:16px;position:relative">
         <span style="position:absolute;left:0;top:2px;width:8px;height:8px;background:COR_HEADER;border-radius:50%;display:inline-block"></span>
         <strong style="color:COR_HEADER;font-size:11px">SUBTÍTULO DO ITEM</strong>
         <p style="margin:3px 0 0;color:#3a3a4a">Explicação completa de 1 a 3 frases descrevendo o mecanismo, procedimento ou conceito de forma educativa e aprofundada.</p>
       </div>
       <div style="margin-bottom:10px;padding-left:16px;position:relative">
         <span style="position:absolute;left:0;top:2px;width:8px;height:8px;background:COR_HEADER;border-radius:50%;display:inline-block"></span>
         <strong style="color:COR_HEADER;font-size:11px">SUBTÍTULO DO ITEM 2</strong>
         <p style="margin:3px 0 0;color:#3a3a4a">Outra explicação completa.</p>
       </div>
     </div>
   </foreignObject>

6. CÁLCULO DE ALTURA:
   - Cada item (subtítulo + parágrafo curto): ~55px; parágrafo médio: ~70px; parágrafo longo: ~90px
   - ALTURA_CORPO = soma dos itens + 24px
   - ALTURA_TOTAL = 46 + ALTURA_CORPO
   - Cards no mesmo par têm a MESMA altura (use o maior)
   - NUNCA deixe foreignObject vazio — todo card tem no mínimo 2 itens explicativos

7. CATEGORIAS E CORES:
   Conceito/Definição: COR_HEADER=#4a6fa5  COR_BG=#eef2f8  COR_BORDA=#c0d0e8
   Fisiopatologia:     COR_HEADER=#7a4040  COR_BG=#faf0f0  COR_BORDA=#d4b0b0
   Sinais e Sintomas:  COR_HEADER=#8a6820  COR_BG=#fdf6e8  COR_BORDA=#ddc878
   Diagnóstico:        COR_HEADER=#2d6a9e  COR_BG=#eaf4fb  COR_BORDA=#a0c4df
   Tratamento:         COR_HEADER=#2d7a50  COR_BG=#edf8f2  COR_BORDA=#90d0b0
   Procedimentos:      COR_HEADER=#5a3a8e  COR_BG=#f2eefa  COR_BORDA=#bba8d8
   Teorias/Modelos:    COR_HEADER=#1a6e6e  COR_BG=#eafafc  COR_BORDA=#88cece
   Cuidados Enfermagem:COR_HEADER=#8e3a60  COR_BG=#faeef4  COR_BORDA=#d8a8bc
   Epidemiologia:      COR_HEADER=#2a6080  COR_BG=#eaf2f8  COR_BORDA=#90b8cc
   Classificação:      COR_HEADER=#5e4a1e  COR_BG=#faf6ee  COR_BORDA=#c8b070

8. Se número ímpar de cards, o último ocupa largura total: x=20, width=920.

9. Retorne APENAS o SVG completo, sem markdown, sem explicação, sem ```"""


class DiagramRequest(BaseModel):
    topic: str
    key_points: list[str]
    sections: list[dict]


@app.post("/api/diagram")
async def generate_diagram(req: DiagramRequest):
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY não configurada.")

    client = anthropic.Anthropic(api_key=api_key)

    content = f"""Tópico: {req.topic}

Pontos-chave:
{chr(10).join(f'- {p}' for p in req.key_points)}

Seções:
{chr(10).join(f'{s["title"]}:{chr(10)}{chr(10).join(f"  - {p}" for p in s["points"])}' for s in req.sections)}

Gere o infográfico SVG com base nesse conteúdo de enfermagem."""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=8192,
            system=SVG_SYSTEM,
            messages=[{"role": "user", "content": content}],
        )
    except anthropic.APIError as e:
        raise HTTPException(status_code=502, detail=str(e))

    svg = response.content[0].text.strip()
    # Remove markdown code fences if Claude added them
    if svg.startswith("```"):
        svg = re.sub(r"^```(?:svg|xml)?\s*", "", svg)
        svg = re.sub(r"\s*```$", "", svg.strip())

    from fastapi.responses import Response
    return Response(content=svg, media_type="image/svg+xml")


class ChatMessage(BaseModel):
    role: Literal["user", "assistant"]
    content: str

class ChatRequest(BaseModel):
    message: str
    context: str
    history: list[ChatMessage] = []


@app.post("/api/chat")
async def chat(req: ChatRequest):
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY não configurada.")

    client = anthropic.Anthropic(api_key=api_key)

    system = [
        {
            "type": "text",
            "text": "Você é um assistente especializado em enfermagem, ajudando uma estudante a entender o conteúdo das aulas.\nResponda de forma clara, didática e objetiva. Use exemplos clínicos quando relevante.\nBaseie suas respostas no seguinte conteúdo estudado:\n\n",
        },
        {
            "type": "text",
            "text": req.context,
            "cache_control": {"type": "ephemeral"},
        },
        {
            "type": "text",
            "text": "\n\nSe a pergunta não estiver relacionada ao conteúdo, responda com base no seu conhecimento geral de enfermagem.",
        },
    ]

    messages = [{"role": m.role, "content": m.content} for m in req.history]
    messages.append({"role": "user", "content": req.message})

    try:
        response = client.messages.create(
            model="claude-haiku-4-5",
            max_tokens=1024,
            system=system,
            messages=messages,
        )
    except anthropic.APIError as e:
        raise HTTPException(status_code=502, detail=f"Erro na API do Claude: {str(e)}")

    return {"reply": response.content[0].text}


@app.get("/api/health")
def health():
    return {"status": "ok"}


# Serve React build in production
static_dir = Path(__file__).parent.parent / "frontend" / "dist"
if static_dir.exists():
    app.mount("/assets", StaticFiles(directory=str(static_dir / "assets")), name="assets")

    @app.get("/{full_path:path}")
    def serve_react(full_path: str):
        candidate = static_dir / full_path
        if candidate.exists() and candidate.is_file():
            return FileResponse(str(candidate))
        return FileResponse(str(static_dir / "index.html"))
